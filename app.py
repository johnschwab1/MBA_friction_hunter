"""
MBA Friction Hunter — Streamlit App

Local dev:
  streamlit run app.py

Deploy to Streamlit Cloud:
  1. Push this repo to GitHub
  2. Connect at share.streamlit.io
  3. Add secrets in the dashboard (see .streamlit/secrets.toml.example)

Provider switching:
  Set LLM_PROVIDER = "anthropic" or "gemini" in your secrets file.

Session persistence:
  Progress is auto-saved to .friction_hunter_session.json after each stage.
  Use the Resume button on the home screen to restore after a page reload.
  Note: this file is local and not committed to git.
"""

import io
import json
import os
import re
from pathlib import Path

import streamlit as st

# ── Bootstrap secrets ──
for _k in ["ANTHROPIC_API_KEY", "ANTHROPIC_MODEL", "GEMINI_API_KEY", "GEMINI_MODEL", "LLM_PROVIDER"]:
    if _k in st.secrets:
        os.environ[_k] = st.secrets[_k]

from llm_client import call_llm, provider_label  # noqa: E402

st.set_page_config(page_title="MBA Friction Hunter", page_icon="🎓", layout="wide", initial_sidebar_state="collapsed")

REPO_ROOT = Path(__file__).parent
SESSION_FILE = REPO_ROOT / ".friction_hunter_session.json"

STAGE_NAMES = {
    2: "Friction Hunter", 3: "Use Case Selector", 4: "Architecture Hunter",
    5: "Design Brief", 6: "Prototype", 7: "Professor Brief", 8: "Reference Document"
}
STAGE_FILES = {
    2: "stage_2_friction_hunter.md", 3: "stage_3_use_case_selector.md",
    4: "stage_4_architecture_hunter.md", 5: "stage_5_design_brief_generator.md",
    6: "stage_6_prototype_generator.md",
    7: "stage_7_professor_brief.md", 8: "stage_8_reference_doc.md",
}
STAGE_SLUGS = {
    2: "friction_hunter", 3: "use_case_selector", 4: "architecture_hunter",
    5: "design_brief", 6: "prototype", 7: "professor_brief", 8: "reference_document"
}
SEARCH_STAGES = {4}
UC_STAGES = (4, 5, 6)

_FAIL_PHRASES = [
    "i notice that", "i cannot", "i'm unable", "i am unable",
    "unable to find", "unable to access", "unable to retrieve",
    "don't have access", "do not have access",
    "no public syllabus", "no publicly available",
    "search results don't contain", "search results do not contain",
    "not find the full text", "couldn't find", "could not find",
    "not publicly available", "not available online",
    "i don't see", "i do not see",
]


def _looks_like_syllabus(text: str) -> bool:
    if len(text) < 500:
        return False
    return not any(p in text[:600].lower() for p in _FAIL_PHRASES)


# ── Session state ──
for _k, _v in {
    "syllabus": "", "outputs": {}, "stage": 0,
    "search_results": "", "search_fetched": False, "search_valid": None,
    "stage3_direction": "", "stage3_examples": "",
    "profile_mode": False, "profile_result": "",
    "use_cases_extracted": [],
    "selected_use_cases": [],
    "use_cases_confirmed": False,
    "auto_running": False,
    "reviewing_prototypes": False,
}.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v


# ── Session persistence ──

def _serialize_outputs(outputs: dict) -> dict:
    result = {}
    for k, v in outputs.items():
        result[str(k)] = {str(uc_k): uc_v for uc_k, uc_v in v.items()} if isinstance(v, dict) else v
    return result


def _deserialize_outputs(data: dict) -> dict:
    result = {}
    for k, v in data.items():
        result[int(k)] = {int(uc_k): uc_v for uc_k, uc_v in v.items()} if isinstance(v, dict) else v
    return result


def save_session():
    try:
        SESSION_FILE.write_text(json.dumps({
            "syllabus": st.session_state.syllabus,
            "outputs": _serialize_outputs(st.session_state.outputs),
            "stage": st.session_state.stage,
            "use_cases_extracted": st.session_state.use_cases_extracted,
            "selected_use_cases": st.session_state.selected_use_cases,
            "use_cases_confirmed": st.session_state.use_cases_confirmed,
            "reviewing_prototypes": st.session_state.reviewing_prototypes,
            "stage3_direction": st.session_state.stage3_direction,
        }, indent=2))
    except Exception:
        pass  # Don't crash the app on save failure


def load_session() -> bool:
    if not SESSION_FILE.exists():
        return False
    try:
        data = json.loads(SESSION_FILE.read_text())
        st.session_state.syllabus = data.get("syllabus", "")
        st.session_state.outputs = _deserialize_outputs(data.get("outputs", {}))
        st.session_state.stage = data.get("stage", 0)
        st.session_state.use_cases_extracted = data.get("use_cases_extracted", [])
        st.session_state.selected_use_cases = data.get("selected_use_cases", [])
        st.session_state.use_cases_confirmed = data.get("use_cases_confirmed", False)
        st.session_state.reviewing_prototypes = data.get("reviewing_prototypes", False)
        st.session_state.stage3_direction = data.get("stage3_direction", "")
        st.session_state.auto_running = False
        return True
    except Exception:
        return False


def clear_session():
    if SESSION_FILE.exists():
        SESSION_FILE.unlink()


# ── Pipeline helpers ──

def get_next_uc_task():
    """Returns (stage, uc_idx) for the next incomplete UC stage, or None."""
    for uc_idx in st.session_state.selected_use_cases:
        for s in UC_STAGES:
            if uc_idx not in st.session_state.outputs.get(s, {}):
                return (s, uc_idx)
    return None


def get_next_report_task():
    """Returns 7 or 8 if that report stage is not yet done, else None."""
    for s in (7, 8):
        if s not in st.session_state.outputs:
            return s
    return None


def get_pipeline_progress():
    """Returns (done, total, status_lines) for the UC stages."""
    selected = st.session_state.selected_use_cases
    uc_list = st.session_state.use_cases_extracted
    total = len(selected) * 3
    done = 0
    lines = []
    for s in UC_STAGES:
        stage_outputs = st.session_state.outputs.get(s, {})
        parts = []
        for uc_idx in selected:
            if uc_idx in stage_outputs:
                parts.append(f"✅ {uc_list[uc_idx]['title']}")
                done += 1
            else:
                parts.append(f"⏳ {uc_list[uc_idx]['title']}")
        lines.append(f"**Stage {s} — {STAGE_NAMES[s]}:** " + " · ".join(parts))
    return done, total, lines


def extract_system_prompt(filename: str) -> str:
    text = (REPO_ROOT / filename).read_text()
    match = re.search(r"## System Prompt\s*\n(.*)", text, re.DOTALL)
    if not match:
        raise ValueError(f"No '## System Prompt' section in {filename}")
    return match.group(1).strip()


@st.cache_data(show_spinner=False)
def load_criteria() -> str:
    return (REPO_ROOT / "stage_1_criteria_layer.md").read_text()


def parse_uploaded_file(f) -> str:
    ext = f.name.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        import pdfplumber
        with pdfplumber.open(io.BytesIO(f.read())) as pdf:
            return "\n\n".join(page.extract_text() or "" for page in pdf.pages)
    if ext == "docx":
        import docx
        doc = docx.Document(io.BytesIO(f.read()))
        return "\n".join(p.text for p in doc.paragraphs)
    if ext == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(f.read()))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    return f.read().decode("utf-8", errors="replace")


def extract_use_cases_from_stage3(stage3_output: str) -> list:
    system = (
        "Extract the distinct use cases from this MBA course AI analysis output. "
        "Return ONLY a valid JSON array. Each element must have exactly two fields: "
        "'title' (the short name of the use case, 1 line) and "
        "'text' (the complete description of that use case as it appears in the document, preserving all detail). "
        "Return raw JSON only — no markdown fences, no explanation, no other text."
    )
    raw = call_llm(system, stage3_output, max_tokens=3000)
    match = re.search(r"\[.*\]", raw, re.DOTALL)
    if match:
        try:
            result = json.loads(match.group())
            if isinstance(result, list) and all("title" in r and "text" in r for r in result):
                return result
        except json.JSONDecodeError:
            pass
    return [{"title": "All Use Cases", "text": stage3_output}]


def build_uc_user_msg(stage: int, uc_idx: int) -> str:
    uc = st.session_state.use_cases_extracted[uc_idx]
    if stage == 4:
        return f"USE CASE TO ANALYZE:\n\n{uc['text']}"
    if stage == 5:
        return (
            f"USE CASE:\n\n{uc['text']}\n\n---\n\n"
            f"ARCHITECTURE RECOMMENDATIONS:\n\n{st.session_state.outputs[4][uc_idx]}"
        )
    if stage == 6:
        return st.session_state.outputs[5][uc_idx]
    return ""


def build_report_user_msg() -> str:
    outputs = st.session_state.outputs
    selected = st.session_state.selected_use_cases
    uc_list = st.session_state.use_cases_extracted
    sections = []
    for uc_idx in selected:
        uc = uc_list[uc_idx]
        sections.append(
            f"=== USE CASE: {uc['title']} ===\n\n"
            f"USE CASE DESCRIPTION:\n{uc['text']}\n\n"
            f"ARCHITECTURE RESEARCH:\n{outputs.get(4, {}).get(uc_idx, '')}\n\n"
            f"DESIGN BRIEF:\n{outputs.get(5, {}).get(uc_idx, '')}\n\n"
            f"PROTOTYPE PROMPT:\n{outputs.get(6, {}).get(uc_idx, '')}"
        )
    return (
        f"FRICTION MAP (Stage 2):\n\n{outputs.get(2, '')}\n\n---\n\n"
        f"ALL SELECTED USE CASES:\n\n" + "\n\n---\n\n".join(sections)
    )


def dict_output_to_md(stage_num: int, output_dict: dict) -> str:
    uc_list = st.session_state.use_cases_extracted
    parts = []
    for uc_idx, text in output_dict.items():
        title = uc_list[uc_idx]["title"] if uc_idx < len(uc_list) else f"Use Case {uc_idx + 1}"
        parts.append(f"# {title}\n\n{text}")
    return "\n\n---\n\n".join(parts)


def _add_runs_with_bold(paragraph, text: str):
    parts = text.split("**")
    for i, part in enumerate(parts):
        if part:
            run = paragraph.add_run(part)
            run.bold = (i % 2 == 1)


def _add_markdown_to_doc(doc, text: str):
    for line in text.split("\n"):
        s = line.rstrip()
        if s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith("- ") or s.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            _add_runs_with_bold(p, s[2:])
        elif re.match(r"^\d+\. ", s):
            p = doc.add_paragraph(style="List Number")
            _add_runs_with_bold(p, re.sub(r"^\d+\. ", "", s))
        elif s == "---":
            doc.add_paragraph("─" * 60)
        elif s:
            p = doc.add_paragraph()
            _add_runs_with_bold(p, s)


def generate_report_docx(professor_brief: str, reference_doc: str) -> bytes:
    from docx import Document
    doc = Document()
    _add_markdown_to_doc(doc, professor_brief)
    doc.add_page_break()
    _add_markdown_to_doc(doc, reference_doc)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def reset():
    clear_session()
    for k, v in {
        "syllabus": "", "outputs": {}, "stage": 0,
        "search_results": "", "search_fetched": False, "search_valid": None,
        "stage3_direction": "", "stage3_examples": "",
        "profile_mode": False, "profile_result": "",
        "use_cases_extracted": [], "selected_use_cases": [], "use_cases_confirmed": False,
        "auto_running": False, "reviewing_prototypes": False,
    }.items():
        st.session_state[k] = v
    st.rerun()


def reset_search():
    st.session_state.search_results = ""
    st.session_state.search_fetched = False
    st.session_state.search_valid = None
    st.session_state.syllabus = ""
    st.session_state.profile_mode = False
    st.session_state.profile_result = ""


# ── Header ──
col_title, col_badge = st.columns([5, 1])
with col_title:
    st.title("🎓 MBA Friction Hunter")
    st.caption("Find where AI creates genuine learning value in an MBA curriculum — then build a working prototype.")
with col_badge:
    st.markdown(f"<br><span style='font-size:0.8em;color:gray'>🤖 {provider_label()}</span>", unsafe_allow_html=True)

st.divider()


# ═══════════════════════════════════════════
# INTAKE
# ═══════════════════════════════════════════

if st.session_state.stage == 0:
    st.subheader("Step 1 — Load the Syllabus")

    # ── Resume saved session ──
    if SESSION_FILE.exists():
        st.info("📂 A previous session was found. Resume it to pick up where you left off, or start fresh.")
        r1, r2 = st.columns([2, 2])
        with r1:
            if st.button("▶️ Resume previous session", type="primary"):
                if load_session():
                    st.rerun()
                else:
                    st.error("Could not load session — starting fresh.")
        with r2:
            if st.button("🗑️ Start fresh"):
                clear_session()
                st.rerun()
        st.divider()

    tab_search, tab_upload, tab_paste = st.tabs(["🔍  Search the web", "📁  Upload a file", "📝  Paste text"])

    with tab_search:
        st.caption("Search for a publicly available syllabus. You'll see a list of options before we fetch the full text.")
        if st.session_state.profile_mode:
            st.info("**Building a course profile from public information.**\n\nEnter the professor's name and course title.")
            prof_name = st.text_input("Professor name", placeholder="e.g. Mathijs de Vaan", key="prof_name_input")
            course_name = st.text_input("Course name / topic", placeholder="e.g. Leading People — Haas MBA", key="course_name_input")
            col_build, col_back = st.columns([2, 2])
            with col_build:
                if st.button("Build profile →", type="primary", disabled=not (prof_name.strip() and course_name.strip())):
                    with st.spinner("Searching…"):
                        st.session_state.profile_result = call_llm(
                            "You are a research assistant. Construct a usable MBA course profile by searching for: "
                            "(1) the professor's research and teaching philosophy; (2) the school's course description; "
                            "(3) typical MBA course structure for this topic at peer institutions. "
                            "Sections: Professor Background, Course Overview & Objectives, Likely Topics & Weekly Structure, "
                            "Key Assignments, Pedagogical Approach. Label inferred content [Inferred].",
                            f"Professor: {prof_name.strip()}\nCourse: {course_name.strip()}",
                            use_search=True, max_tokens=3000
                        )
                    st.rerun()
            with col_back:
                if st.button("← Back"):
                    reset_search(); st.rerun()
            if st.session_state.profile_result:
                st.success(f"✓ Profile built — {len(st.session_state.profile_result):,} chars")
                with st.expander("Preview", expanded=True):
                    st.markdown(st.session_state.profile_result[:2000] +
                                ("\n\n*[truncated]*" if len(st.session_state.profile_result) > 2000 else ""))
                st.caption("⚠️ Synthesized profile, not the actual syllabus.")
                if st.button("Start analysis →", type="primary"):
                    st.session_state.syllabus = st.session_state.profile_result
                    st.session_state.stage = 2
                    save_session()
                    st.rerun()
        else:
            q_col, btn_col = st.columns([4, 1])
            with q_col:
                query = st.text_input("Query", placeholder="e.g. Haas MBA marketing 2025",
                                      label_visibility="collapsed", key="search_query")
            with btn_col:
                do_search = st.button("Search →", use_container_width=True)
            if do_search and query:
                with st.spinner(f"Searching for '{query}'…"):
                    st.session_state.search_results = call_llm(
                        "You are a research assistant. Search for MBA course syllabi matching the user's query. "
                        "Return a numbered list of up to 5 options. For each: number, title, school, year, "
                        "and one-sentence description. Format: N. Title — School (Year): Description",
                        f"Find MBA syllabi for: {query}", use_search=True, max_tokens=800
                    )
                    st.session_state.search_fetched = False
                    st.session_state.search_valid = None
            if st.session_state.search_results:
                st.markdown(st.session_state.search_results)
                st.divider()
                if st.session_state.search_fetched and st.session_state.search_valid is False:
                    st.warning("**No public syllabus found.** Choose how to proceed:")
                    with st.expander("See what was returned"):
                        st.text(st.session_state.syllabus[:800] + "…")
                    r_col, p_col, m_col = st.columns(3)
                    with r_col:
                        if st.button("🔄 New search query", use_container_width=True):
                            reset_search(); st.rerun()
                    with p_col:
                        if st.button("🧠 Build from public info", use_container_width=True):
                            st.session_state.profile_mode = True; st.rerun()
                    with m_col:
                        st.caption("Or switch to Upload / Paste tab above.")
                elif not st.session_state.search_fetched:
                    pick_col, fetch_col = st.columns([1, 3])
                    with pick_col:
                        pick = st.number_input("Which result?", min_value=1, max_value=5, step=1)
                    with fetch_col:
                        if st.button("Fetch full syllabus →", use_container_width=True):
                            with st.spinner("Fetching…"):
                                fetched = call_llm(
                                    "Fetch the full syllabus text of the selected result. "
                                    "Return only the content: description, objectives, schedule, assignments, readings.",
                                    f"Results:\n\n{st.session_state.search_results}\n\nFetch option {pick}.",
                                    use_search=True, max_tokens=4096
                                )
                            st.session_state.syllabus = fetched
                            st.session_state.search_fetched = True
                            st.session_state.search_valid = _looks_like_syllabus(fetched)
                            st.rerun()
                elif st.session_state.search_valid:
                    st.success(f"✓ Fetched {len(st.session_state.syllabus):,} chars")
                    with st.expander("Preview"):
                        st.text(st.session_state.syllabus[:1000] + "…")
                    c1, c2 = st.columns([2, 2])
                    with c1:
                        if st.button("Start analysis →", type="primary", key="search_go"):
                            st.session_state.stage = 2; save_session(); st.rerun()
                    with c2:
                        if st.button("🔄 Try a different result"):
                            st.session_state.search_fetched = False
                            st.session_state.search_valid = None
                            st.session_state.syllabus = ""
                            st.rerun()

    with tab_upload:
        st.caption("PDF, DOCX, PPTX, TXT. For Google Docs: export as PDF first.")
        uploaded = st.file_uploader("File", type=["pdf", "docx", "pptx", "txt"], label_visibility="collapsed")
        if uploaded:
            with st.spinner("Parsing…"):
                parsed = parse_uploaded_file(uploaded)
            st.success(f"✓ {uploaded.name} — {len(parsed):,} chars")
            with st.expander("Preview"):
                st.text(parsed[:1000] + ("..." if len(parsed) > 1000 else ""))
            if st.button("Start analysis →", type="primary", key="upload_go"):
                st.session_state.syllabus = parsed
                st.session_state.stage = 2
                save_session()
                st.rerun()

    with tab_paste:
        pasted = st.text_area("Text", height=320, placeholder="Paste your syllabus here…", label_visibility="collapsed")
        if st.button("Start analysis →", type="primary", key="paste_go") and pasted.strip():
            st.session_state.syllabus = pasted.strip()
            st.session_state.stage = 2
            save_session()
            st.rerun()


# ═══════════════════════════════════════════
# PIPELINE
# ═══════════════════════════════════════════

elif st.session_state.stage >= 2:
    criteria = load_criteria()
    outputs = st.session_state.outputs

    if st.button("↩ New syllabus"):
        reset()

    # ── STAGE 2 ──
    if st.session_state.stage == 2:
        st.subheader("Stage 2: Friction Hunter")
        if 2 not in outputs:
            if st.button("Run Stage 2 →", type="primary"):
                system = extract_system_prompt(STAGE_FILES[2])
                user_msg = f"CRITERIA REFERENCE:\n\n{criteria}\n\n---\n\nCURRICULUM TO ANALYZE:\n\n{st.session_state.syllabus}"
                with st.spinner("Friction Hunter…"):
                    outputs[2] = call_llm(system, user_msg)
                save_session()
                st.rerun()
        else:
            st.markdown(outputs[2])
            st.divider()
            a_col, dl_col, _ = st.columns([2, 2, 3])
            with a_col:
                if st.button("↺ Re-run"):
                    del outputs[2]; save_session(); st.rerun()
            with dl_col:
                st.download_button("Download .md", data=outputs[2],
                    file_name="stage_02_friction_hunter.md", mime="text/markdown", key="dl2")
            if st.button("Continue to Stage 3 →", type="primary"):
                st.session_state.stage = 3; save_session(); st.rerun()

    # ── STAGE 3 ──
    elif st.session_state.stage == 3:
        with st.expander("✅ Stage 2: Friction Hunter", expanded=False):
            st.markdown(outputs.get(2, ""))
            st.download_button("Download .md", data=outputs.get(2, ""),
                file_name="stage_02_friction_hunter.md", mime="text/markdown", key="dl2_done")

        st.subheader("Stage 3: Use Case Selector")
        if 3 not in outputs:
            if st.button("Run Stage 3 →", type="primary"):
                system = extract_system_prompt(STAGE_FILES[3])
                base = outputs.get(2, st.session_state.syllabus)
                d = st.session_state.stage3_direction
                user_msg = f"{base}\n\nDIRECTION FROM PROFESSOR:\n{d}" if d else base
                with st.spinner("Use Case Selector…"):
                    outputs[3] = call_llm(system, user_msg)
                save_session()
                st.rerun()
        else:
            st.markdown(outputs[3])
            st.divider()

            if st.session_state.use_cases_extracted and not st.session_state.use_cases_confirmed:
                # ── Checkbox selection screen ──
                st.subheader("☑️ Choose which use cases to develop")
                st.caption("Select any or all. Each will go through Architecture → Design Brief → Prototype automatically.")
                selected = []
                for i, uc in enumerate(st.session_state.use_cases_extracted):
                    if st.checkbox(uc["title"], value=True, key=f"uc_check_{i}"):
                        selected.append(i)
                st.divider()
                c1, c2 = st.columns([2, 3])
                with c1:
                    if st.button(f"Run pipeline for {len(selected)} use case(s) →",
                                 type="primary", disabled=len(selected) == 0):
                        st.session_state.selected_use_cases = selected
                        st.session_state.use_cases_confirmed = True
                        st.session_state.auto_running = True
                        st.session_state.stage = 4
                        save_session()
                        st.rerun()
                with c2:
                    if st.button("↺ Re-run Stage 3 with new direction"):
                        st.session_state.use_cases_extracted = []
                        del outputs[3]
                        st.session_state.stage3_examples = ""
                        save_session()
                        st.rerun()

            elif not st.session_state.use_cases_extracted:
                # ── Approval loop ──
                if not st.session_state.stage3_examples:
                    with st.spinner("Generating guidance suggestions…"):
                        st.session_state.stage3_examples = call_llm(
                            "Based on these MBA course AI use case recommendations, write 3-4 short, specific "
                            "examples of direction a professor might give to improve or redirect the selection. "
                            "Make each concrete and tied to the actual use cases shown. "
                            "Format as a markdown bulleted list. Each is one actionable sentence.",
                            f"Use cases:\n\n{outputs[3]}", max_tokens=400
                        )
                st.info("💡 **Not sure what to change?**\n\n" + st.session_state.stage3_examples)
                direction_input = st.text_area(
                    "Give direction and I'll re-run — or leave blank and click Satisfied:",
                    placeholder="e.g. 'The negotiation sim in week 7 is the hardest part — center the use cases there'",
                    key="stage3_direction_input", height=100,
                )
                col1, col2, col3 = st.columns([2, 2, 2])
                with col1:
                    if st.button("↺ Re-run with this direction", disabled=not direction_input.strip()):
                        st.session_state.stage3_direction = direction_input.strip()
                        st.session_state.stage3_examples = ""
                        del outputs[3]
                        save_session()
                        st.rerun()
                with col2:
                    if st.button("✓ Satisfied — Choose use cases →", type="primary"):
                        with st.spinner("Extracting use cases…"):
                            st.session_state.use_cases_extracted = extract_use_cases_from_stage3(outputs[3])
                        st.session_state.stage3_direction = ""
                        save_session()
                        st.rerun()
                with col3:
                    st.download_button("Download .md", data=outputs[3],
                        file_name="stage_03_use_case_selector.md", mime="text/markdown", key="dl3")

    # ── STAGES 4+ ──
    elif st.session_state.stage >= 4:
        uc_list = st.session_state.use_cases_extracted
        selected = st.session_state.selected_use_cases

        # Completed stages 2 & 3 always shown collapsed
        for s in (2, 3):
            if s in outputs:
                with st.expander(f"✅ Stage {s}: {STAGE_NAMES[s]}", expanded=False):
                    st.markdown(outputs[s])
                    st.download_button("Download .md", data=outputs[s],
                        file_name=f"stage_{s:02d}_{STAGE_SLUGS[s]}.md", mime="text/markdown",
                        key=f"dl_done_{s}")

        # ═══════════════════════════════════════════
        # AUTO-RUN: Stages 4-6 per use case
        # ═══════════════════════════════════════════
        if st.session_state.auto_running and not st.session_state.reviewing_prototypes:
            next_uc = get_next_uc_task()
            next_report = get_next_report_task()

            if next_uc:
                s, uc_idx = next_uc
                done, total, lines = get_pipeline_progress()
                st.subheader(f"🔄 Running pipeline — {done}/{total} complete")
                for line in lines:
                    st.markdown(line)
                st.divider()
                uc = uc_list[uc_idx]
                with st.spinner(f"Stage {s} — {STAGE_NAMES[s]}: {uc['title']}…"):
                    system = extract_system_prompt(STAGE_FILES[s])
                    result = call_llm(system, build_uc_user_msg(s, uc_idx),
                                      use_search=(s in SEARCH_STAGES))
                if s not in outputs:
                    outputs[s] = {}
                outputs[s][uc_idx] = result
                save_session()
                st.rerun()

            elif next_report:
                # Report stages
                st.subheader(f"🔄 Generating {STAGE_NAMES[next_report]}…")
                with st.spinner(f"{STAGE_NAMES[next_report]}…"):
                    system = extract_system_prompt(STAGE_FILES[next_report])
                    outputs[next_report] = call_llm(system, build_report_user_msg())
                save_session()
                st.rerun()

            else:
                # Everything done
                st.session_state.auto_running = False
                save_session()
                st.rerun()

        # ═══════════════════════════════════════════
        # REVIEW: Prototype check before report generation
        # ═══════════════════════════════════════════
        elif st.session_state.reviewing_prototypes:
            st.subheader("📋 Review: Prototype Prompts")
            st.caption("The pipeline has finished. Review the prototype prompts below — re-run any that need adjustment — then generate the full report.")

            proto_outputs = outputs.get(6, {})
            if selected and all(i in proto_outputs for i in selected):
                tabs = st.tabs([uc_list[i]["title"] for i in selected])
                for tab, uc_idx in zip(tabs, selected):
                    with tab:
                        st.markdown(proto_outputs[uc_idx])
                        st.divider()
                        r_col, dl_col, _ = st.columns([2, 2, 3])
                        with r_col:
                            if st.button("↺ Re-run this prototype", key=f"rerun_proto_{uc_idx}"):
                                outputs[6].pop(uc_idx, None)
                                st.session_state.auto_running = True
                                save_session()
                                st.rerun()
                        with dl_col:
                            st.download_button("Download .md", data=proto_outputs[uc_idx],
                                file_name=f"prototype_{uc_idx}.md", mime="text/markdown",
                                key=f"dl_proto_{uc_idx}")

            st.divider()
            if st.button("📄 Generate Full Report →", type="primary"):
                st.session_state.reviewing_prototypes = False
                st.session_state.stage = 7
                st.session_state.auto_running = True
                save_session()
                st.rerun()

        # ═══════════════════════════════════════════
        # COMPLETE: Downloads
        # ═══════════════════════════════════════════
        elif 7 in outputs and 8 in outputs:
            st.success("🏁 Pipeline complete!")
            st.subheader("📥 Download Report")
            docx_bytes = generate_report_docx(outputs[7], outputs[8])
            dl1, dl2, dl3 = st.columns(3)
            with dl1:
                st.download_button(
                    label="📄 Full Report (.docx)",
                    data=docx_bytes,
                    file_name="mba_friction_hunter_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    type="primary"
                )
            with dl2:
                st.download_button("Professor Brief (.md)", data=outputs[7],
                    file_name="professor_brief.md", mime="text/markdown")
            with dl3:
                st.download_button("Reference Doc (.md)", data=outputs[8],
                    file_name="reference_document.md", mime="text/markdown")
            st.caption("💡 To open in Google Docs: upload the .docx to Google Drive → right-click → Open with Google Docs.")
            st.divider()

            # Re-run report stages if needed
            rr1, rr2 = st.columns([2, 2])
            with rr1:
                if st.button("↺ Re-run Professor Brief"):
                    del outputs[7]
                    st.session_state.reviewing_prototypes = False
                    st.session_state.stage = 7
                    st.session_state.auto_running = True
                    save_session(); st.rerun()
            with rr2:
                if st.button("↺ Re-run Reference Doc"):
                    del outputs[8]
                    st.session_state.reviewing_prototypes = False
                    st.session_state.stage = 7
                    st.session_state.auto_running = True
                    save_session(); st.rerun()

            st.divider()
            st.subheader("All Stage Outputs")
            for s in range(2, 9):
                if s not in outputs:
                    continue
                out = outputs[s]
                dl_data = dict_output_to_md(s, out) if isinstance(out, dict) else out
                st.download_button(
                    label=f"Stage {s}: {STAGE_NAMES[s]}",
                    data=dl_data,
                    file_name=f"stage_{s:02d}_{STAGE_SLUGS[s]}.md",
                    mime="text/markdown",
                    key=f"dl_final_{s}"
                )

        # ── Fallback: if auto_run finished UC stages, trigger review ──
        elif not st.session_state.auto_running and not st.session_state.reviewing_prototypes:
            uc_done = get_next_uc_task() is None and len(selected) > 0
            report_needed = 7 not in outputs
            if uc_done and report_needed:
                st.session_state.reviewing_prototypes = True
                save_session()
                st.rerun()
